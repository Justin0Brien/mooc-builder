#!/usr/bin/env python3
"""
mooc_builder.py

Scan a directory tree of legacy teaching materials (PowerPoint, Keynote, Word, PDF,
plain text, images, and videos), extract their text, and use a local LLM via
`llama-cpp-python` (the Llama class) to:
  1) normalise and rewrite content into clean Markdown suitable for MOOCs;
  2) auto-classify each source into a logical subject folder
     (Statistics, Research Methods, Neuroimaging, Vision Science, Personality,
      Individual Differences, or Other/Discovered);
  3) generate helpful extras (lecture notes, slide outline, quiz items).

Outputs a new folder structure containing ONLY Markdown files (plus a small cache and an index file).

Designed for macOS and Linux. Keynote support on macOS uses AppleScript to export to PDF.

Requirements (install as needed):
  pip install --upgrade llama-cpp-python pdfminer.six python-docx python-pptx pytesseract pillow tqdm pyyaml
Optional for video:  pip install openai-whisper
System deps (optional but helpful):
  - Tesseract OCR (for images): macOS (brew install tesseract), Ubuntu (apt-get install tesseract-ocr)
  - ffmpeg (for extracting audio from video): brew install ffmpeg / apt-get install ffmpeg

Example:
  python mooc_builder.py \
    --input-dir ~/Documents/Teaching \
    --output-dir ~/MOOC_Build \
    --model /path/to/model.gguf \
    --ctx 8192 --max-tokens 1536 --temperature 0.2

Notes:
- This script keeps all processing local. It never uploads your data.
- It is conservative about Keynote, images, and video: if tools are missing, it skips gracefully with a warning.
- British English spelling is enforced in prompts.
"""
from __future__ import annotations

import argparse
import concurrent.futures
import hashlib
import json
import logging
import os
import re
import shutil
import subprocess
import sys
import tempfile
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from typing import Dict, Iterable, List, Optional, Tuple

# Third-party
# Switched from llama-cpp-python to Ollama local server API.
try:  # ollama python client
    import ollama  # type: ignore
except Exception:
    ollama = None  # type: ignore

try:
    from pdfminer.high_level import extract_text as pdf_extract_text
except Exception:
    pdf_extract_text = None

try:
    import docx  # python-docx
except Exception:
    docx = None

try:
    from pptx import Presentation  # python-pptx
except Exception:
    Presentation = None

try:
    from PIL import Image
    import pytesseract
except Exception:
    Image = None  # type: ignore
    pytesseract = None  # type: ignore

# Whisper is optional
try:
    import whisper  # openai-whisper
except Exception:
    whisper = None

try:
    import yaml
except Exception:
    yaml = None

from tqdm import tqdm

SUPPORTED_TEXT_EXTS = {".txt", ".md", ".rst"}
SUPPORTED_PDF_EXTS = {".pdf"}
SUPPORTED_DOCX_EXTS = {".docx"}
SUPPORTED_PPTX_EXTS = {".pptx"}
SUPPORTED_KEYNOTE_EXTS = {".key"}
SUPPORTED_IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".tiff", ".bmp", ".gif"}
SUPPORTED_VIDEO_EXTS = {".mp4", ".mov", ".m4v", ".mkv", ".avi"}

DEFAULT_SUBJECTS = [
    "Statistics",
    "Research Methods",
    "Neuroimaging",
    "Vision Science",
    "Personality",
    "Individual Differences",
    "Other",
]

# --- Prompts ------------------------------------------------------------------
SYSTEM_PROMPT = (
    "You are an expert UK psychology lecturer and curriculum designer. "
    "You write in clear British English and avoid Americanisms. "
    "You will receive raw extracted content from legacy teaching files (slides, PDFs, docs). "
    "Your task is to produce a single, self-contained Markdown document suitable for a MOOC. "
    "Do NOT add facts that are not supported by the provided content; no hallucinations. "
    "If content is thin or fragmentary, keep it concise and mark gaps with TODOs."
)

MD_INSTRUCTIONS = (
    "Rewrite and normalise the provided content into Markdown with this structure:\n\n"
    "# Title\n"
    "(Pick a short, apt title. If unclear, infer from content conservatively.)\n\n"
    "## Overview\n"
    "A brief summary in 3–6 sentences, in British English.\n\n"
    "## Key Learning Outcomes\n"
    "- Outcome 1\n- Outcome 2\n- Outcome 3\n"
    "(Replace with specific outcomes based on the content.)\n\n"
    "## Lecture Notes\n"
    "Well-organised narrative covering the material. Use subsections (###) as needed.\n"
    "Cite embedded figures/tables from the original if present, e.g., 'Figure: ...' (no actual images).\n\n"
    "## Slide Outline\n"
    "Bullet-point outline for slides (10–20 bullets).\n\n"
    "## Examples and Activities\n"
    "Short exercises or discussion prompts drawn from the material.\n\n"
    "## Quiz\n"
    "Provide 5–10 multiple-choice questions (A–D) with answers and one-line explanations.\n\n"
    "## References (as extracted)\n"
    "List any references explicitly present in the source. If none, say 'None given'.\n\n"
    "## Notes\n"
    "Mention any uncertainties or TODOs for the instructor to review."
)

CATEGORY_PROMPT = (
    "Classify the provided content into ONE of these subject areas, "
    "or propose a single sensible alternative if none fit well:\n"
    + "\n".join(f"- {s}" for s in DEFAULT_SUBJECTS[:-1]) +
    "\n- Other\n\n"
    "Reply ONLY with the category name."
)

@dataclass
class Config:
    input_dir: Path
    output_dir: Path
    model_name: str  # Ollama model name (e.g., gpt-oss, llama3.2, etc.)
    ctx_size: int = 8192
    threads: int = max(2, os.cpu_count() or 2)  # retained (may be unused by Ollama)
    temperature: float = 0.2
    max_tokens: int = 1536
    ocr: bool = False
    transcribe_video: bool = False
    whisper_model: str = "small"  # used if transcribe_video
    dry_run: bool = False
    ignore_globs: List[str] = None  # type: ignore


# --- Utilities ----------------------------------------------------------------

def sha1_of_path(p: Path) -> str:
    h = hashlib.sha1()
    h.update(str(p).encode("utf-8"))
    try:
        h.update(str(p.stat().st_mtime_ns).encode("utf-8"))
        h.update(str(p.stat().st_size).encode("utf-8"))
    except Exception:
        pass
    return h.hexdigest()


def load_cache(cache_path: Path) -> Dict[str, dict]:
    if cache_path.exists():
        try:
            return json.loads(cache_path.read_text())
        except Exception:
            logging.warning("Cache file corrupted; starting fresh.")
    return {}


def save_cache(cache_path: Path, cache: Dict[str, dict]) -> None:
    tmp = cache_path.with_suffix(".tmp")
    tmp.write_text(json.dumps(cache, indent=2))
    tmp.replace(cache_path)


# --- Extraction ---------------------------------------------------------------

def extract_text_from_txt(path: Path) -> str:
    return path.read_text(errors="ignore")


def extract_text_from_pdf(path: Path) -> str:
    if pdf_extract_text is None:
        raise RuntimeError("pdfminer.six not installed")
    return pdf_extract_text(str(path))


def extract_text_from_docx(path: Path) -> str:
    if docx is None:
        raise RuntimeError("python-docx not installed")
    d = docx.Document(str(path))
    parts = [p.text for p in d.paragraphs]
    # tables
    for table in d.tables:
        for row in table.rows:
            parts.append("\t".join(cell.text for cell in row.cells))
    return "\n".join(parts)


def extract_text_from_pptx(path: Path) -> str:
    if Presentation is None:
        raise RuntimeError("python-pptx not installed")
    prs = Presentation(str(path))
    texts = []
    for i, slide in enumerate(prs.slides, start=1):
        texts.append(f"\n--- Slide {i} ---\n")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
    return "\n".join(texts)


def extract_text_from_keynote(path: Path) -> str:
    # macOS only: export to temp PDF via AppleScript then parse
    if sys.platform != "darwin":
        raise RuntimeError("Keynote extraction only supported on macOS")
    tmpdir = Path(tempfile.mkdtemp(prefix="keynote_export_"))
    pdf_path = tmpdir / (path.stem + ".pdf")
    ascript = f'''
    tell application "Keynote"
        open POSIX file "{path.as_posix()}"
        delay 0.5
        export front document to POSIX file "{pdf_path.as_posix()}" as PDF
        close front document saving no
    end tell
    '''
    try:
        res = subprocess.run(["osascript", "-e", ascript], capture_output=True, text=True, timeout=300)
        if res.returncode != 0:
            raise RuntimeError(f"AppleScript export failed: {res.stderr.strip()}")
        text = extract_text_from_pdf(pdf_path)
        return text
    finally:
        try:
            shutil.rmtree(tmpdir)
        except Exception:
            pass


def extract_text_from_image(path: Path, ocr: bool) -> str:
    if not ocr:
        return ""  # OCR disabled
    if Image is None or pytesseract is None:
        raise RuntimeError("pytesseract/Pillow not installed; cannot OCR images")
    with Image.open(str(path)) as im:
        im = im.convert("L")  # greyscale tends to help OCR
        return pytesseract.image_to_string(im)


def extract_text_from_video(path: Path, transcribe: bool, whisper_model: str) -> str:
    if not transcribe:
        return ""
    if whisper is None:
        raise RuntimeError("openai-whisper not installed; cannot transcribe video")
    # Extract audio to temp wav via ffmpeg
    tmpdir = Path(tempfile.mkdtemp(prefix="vid_audio_"))
    wav_path = tmpdir / (path.stem + ".wav")
    try:
        ff = subprocess.run([
            "ffmpeg", "-y", "-i", str(path),
            "-vn", "-ac", "1", "-ar", "16000", str(wav_path)
        ], capture_output=True)
        if ff.returncode != 0:
            raise RuntimeError("ffmpeg failed to extract audio")
        model = whisper.load_model(whisper_model)
        result = model.transcribe(str(wav_path))
        return result.get("text", "").strip()
    finally:
        try:
            shutil.rmtree(tmpdir)
        except Exception:
            pass


def extract_text(path: Path, cfg: Config) -> str:
    ext = path.suffix.lower()
    try:
        if ext in SUPPORTED_TEXT_EXTS:
            return extract_text_from_txt(path)
        if ext in SUPPORTED_PDF_EXTS:
            return extract_text_from_pdf(path)
        if ext in SUPPORTED_DOCX_EXTS:
            return extract_text_from_docx(path)
        if ext in SUPPORTED_PPTX_EXTS:
            return extract_text_from_pptx(path)
        if ext in SUPPORTED_KEYNOTE_EXTS:
            return extract_text_from_keynote(path)
        if ext in SUPPORTED_IMAGE_EXTS:
            return extract_text_from_image(path, cfg.ocr)
        if ext in SUPPORTED_VIDEO_EXTS:
            return extract_text_from_video(path, cfg.transcribe_video, cfg.whisper_model)
    except Exception as e:
        logging.warning(f"Failed to extract {path.name}: {e}")
        return ""
    return ""  # unsupported


# --- LLM calls (Ollama) -------------------------------------------------------

def ensure_ollama_available():
    if ollama is None:
        raise RuntimeError("ollama python package not installed. pip install ollama")


def chat_completion(model_name: str, system_prompt: str, user_prompt: str, cfg: Config) -> str:
    """Call local Ollama server using chat endpoint.

    Parameters mapped:
      temperature -> options.temperature
      ctx_size -> options.num_ctx
      max_tokens -> options.num_predict
    """
    ensure_ollama_available()
    messages = []
    if system_prompt:
        messages.append({"role": "system", "content": system_prompt})
    messages.append({"role": "user", "content": user_prompt})
    try:
        res = ollama.chat(
            model=model_name,
            messages=messages,
            options={
                "temperature": cfg.temperature,
                "num_ctx": cfg.ctx_size,
                "num_predict": cfg.max_tokens,
            },
            stream=False,
        )
        return res.get("message", {}).get("content", "").strip()
    except Exception as e:
        logging.warning(f"Ollama chat failed: {e}")
        return ""


def approximate_chunks(text: str, target_chars: int = 12000) -> List[str]:
    text = text.strip()
    if len(text) <= target_chars:
        return [text]
    # Split on paragraph boundaries near target size
    paras = re.split(r"\n\n+", text)
    chunks = []
    current = []
    size = 0
    for p in paras:
        if size + len(p) > target_chars and current:
            chunks.append("\n\n".join(current))
            current = [p]
            size = len(p)
        else:
            current.append(p)
            size += len(p)
    if current:
        chunks.append("\n\n".join(current))
    return chunks


def llm_markdown_from_text(model_name: str, raw_text: str, cfg: Config) -> str:
    final_sections = []
    chunks = approximate_chunks(raw_text)
    for i, chunk in enumerate(chunks, start=1):
        user_prompt = (
            "Source chunk " + str(i) + "/" + str(len(chunks)) + " (verbatim, may be messy):\n\n" + chunk +
            "\n\n---\n\nNow, " + MD_INSTRUCTIONS
        )
        md = chat_completion(model_name, SYSTEM_PROMPT, user_prompt, cfg)
        if not md:
            continue
        # For multi-chunk sources, annotate section
        if len(chunks) > 1:
            md = f"\n\n<!-- Chunk {i} of {len(chunks)} -->\n\n" + md
        final_sections.append(md)
    if not final_sections:
        return ""
    # If multiple, we could attempt a consolidation pass
    if len(final_sections) == 1:
        return final_sections[0]
    else:
        combined = "\n\n".join(final_sections)
        # Optional consolidation: ask LLM to merge sections into one coherent doc
        merge_prompt = (
            "Combine the following Markdown sections (labelled by HTML comments) into a single, coherent Markdown document. "
            "Keep the requested headings. Remove duplicates. British English only.\n\n" + combined
        )
        merged = chat_completion(model_name, SYSTEM_PROMPT, merge_prompt, cfg)
        return merged or combined


def llm_categorise(model_name: str, text: str, cfg: Config) -> str:
    snip = text[:4000]
    cat = chat_completion(model_name, SYSTEM_PROMPT, f"{CATEGORY_PROMPT}\n\nCONTENT:\n\n{snip}", cfg)
    if not cat:
        return "Other"
    cat = cat.strip().splitlines()[0]
    # Normalise to one of DEFAULT_SUBJECTS or accept a new single token category
    norm = {
        s.lower(): s for s in DEFAULT_SUBJECTS
    }
    ckey = cat.lower().strip()
    if ckey in norm:
        return norm[ckey]
    # Attempt fuzzy match
    for s in DEFAULT_SUBJECTS:
        if s.lower() in ckey or ckey in s.lower():
            return s
    # Keep a conservative simple result
    return re.sub(r"[^A-Za-z0-9 &-]", "", cat)[:60] or "Other"


# --- Orchestration ------------------------------------------------------------

def should_ignore(path: Path, ignore_globs: Optional[List[str]]) -> bool:
    if not ignore_globs:
        return False
    for pat in ignore_globs:
        if path.match(pat):
            return True
    return False


def discover_files(root: Path, ignore_globs: Optional[List[str]]) -> List[Path]:
    files = []
    for p in root.rglob("*"):
        if p.is_file() and not should_ignore(p, ignore_globs):
            ext = p.suffix.lower()
            if (
                ext in SUPPORTED_TEXT_EXTS
                or ext in SUPPORTED_PDF_EXTS
                or ext in SUPPORTED_DOCX_EXTS
                or ext in SUPPORTED_PPTX_EXTS
                or ext in SUPPORTED_KEYNOTE_EXTS
                or ext in SUPPORTED_IMAGE_EXTS
                or ext in SUPPORTED_VIDEO_EXTS
            ):
                files.append(p)
    return files


def write_markdown(output_dir: Path, category: str, source_path: Path, markdown: str) -> Path:
    # Safe category folder name
    cat_folder = re.sub(r"[^A-Za-z0-9._ -]", "", category).strip() or "Other"
    out_dir = output_dir / cat_folder
    out_dir.mkdir(parents=True, exist_ok=True)

    # File name based on source
    base_name = re.sub(r"\s+", "_", source_path.stem)
    base_name = re.sub(r"[^A-Za-z0-9._-]", "", base_name)
    out_path = out_dir / f"{base_name}.md"

    # Front matter
    front = {
        "title": source_path.stem,
        "category": category,
        "source_path": str(source_path),
        "processed_at": datetime.now().isoformat(timespec="seconds"),
        "tags": [category],
        "model": "ollama",
        "language": "en-GB",
    }
    fm = "---\n" + (yaml.safe_dump(front) if yaml else json.dumps(front, indent=2)) + "---\n\n"

    out_path.write_text(fm + markdown)
    return out_path


def build_index(output_dir: Path) -> None:
    entries = []
    for md in output_dir.rglob("*.md"):
        if md.name == "INDEX.md":
            continue
        rel = md.relative_to(output_dir)
        parts = rel.parts
        category = parts[0] if len(parts) > 1 else "Other"
        title = md.stem.replace("_", " ")
        entries.append((category, str(rel), title))
    entries.sort()
    lines = ["# MOOC Build Index\n"]
    current_cat = None
    for cat, rel, title in entries:
        if cat != current_cat:
            lines.append(f"\n## {cat}\n")
            current_cat = cat
        lines.append(f"- [{title}]({rel})")
    (output_dir / "INDEX.md").write_text("\n".join(lines) + "\n")


# --- Main ---------------------------------------------------------------------

def process_one(path: Path, cfg: Config, model_name: str, cache: Dict[str, dict]) -> Optional[Path]:
    key = sha1_of_path(path)
    if key in cache:
        return Path(cache[key]["out"]) if cache[key].get("out") else None

    raw = extract_text(path, cfg)
    if not raw or not raw.strip():
        logging.info(f"No extractable text in {path}")
        cache[key] = {"out": None, "skip": True}
        return None

    if cfg.dry_run:
        cache[key] = {"out": None, "skip": False}
        return None

    md = llm_markdown_from_text(model_name, raw, cfg)
    if not md:
        logging.warning(f"LLM produced no output for {path}")
        cache[key] = {"out": None, "skip": True}
        return None

    category = llm_categorise(model_name, md, cfg)
    out_path = write_markdown(cfg.output_dir, category, path, md)
    cache[key] = {"out": str(out_path), "skip": False}
    return out_path


def main():
    ap = argparse.ArgumentParser(description="Build a Markdown MOOC from legacy teaching files using a local LLM")
    ap.add_argument("--input-dir", required=True, type=Path)
    ap.add_argument("--output-dir", required=True, type=Path)
    ap.add_argument("--model", required=True, type=str, help="Ollama model name (e.g. gpt-oss, llama3.2, mistral, etc.)")
    ap.add_argument("--ctx", type=int, default=8192, help="Context window size to request (num_ctx)")
    ap.add_argument("--threads", type=int, default=max(2, os.cpu_count() or 2))
    ap.add_argument("--temperature", type=float, default=0.2)
    ap.add_argument("--max-tokens", type=int, default=1536)
    ap.add_argument("--ocr", action="store_true", help="Enable OCR for images via Tesseract")
    ap.add_argument("--transcribe-video", action="store_true", help="Transcribe video audio with Whisper (requires ffmpeg and whisper)")
    ap.add_argument("--whisper-model", type=str, default="small")
    ap.add_argument("--ignore", action="append", default=[], help="Glob(s) to ignore (can repeat)")
    ap.add_argument("--dry-run", action="store_true", help="Scan and extract but do not call LLM or write outputs")
    args = ap.parse_args()

    logging.basicConfig(level=logging.INFO, format="[%(levelname)s] %(message)s")

    cfg = Config(
        input_dir=args.input_dir.expanduser().resolve(),
        output_dir=args.output_dir.expanduser().resolve(),
    model_name=args.model.strip(),
        ctx_size=args.ctx,
        threads=args.threads,
        temperature=args.temperature,
        max_tokens=args.max_tokens,
        ocr=args.ocr,
        transcribe_video=args.transcribe_video,
        whisper_model=args.whisper_model,
        dry_run=args.dry_run,
        ignore_globs=args.ignore,
    )

    cfg.output_dir.mkdir(parents=True, exist_ok=True)
    cache_path = cfg.output_dir / ".mooc_cache.json"
    cache = load_cache(cache_path)

    files = discover_files(cfg.input_dir, cfg.ignore_globs)
    if not files:
        logging.error("No supported files found.")
        sys.exit(1)

    logging.info(f"Discovered {len(files)} files to consider.")

    model_name = cfg.model_name
    if not cfg.dry_run:
        # Warm load: empty chat to ensure model loads (ignore failures)
        try:
            chat_completion(model_name, "", "", cfg)
            logging.info(f"Ollama model '{model_name}' ready.")
        except Exception as e:
            logging.error(f"Failed to load Ollama model '{model_name}': {e}")
            sys.exit(1)

    try:
        for p in tqdm(files, desc="Processing"):
            try:
                process_one(p, cfg, model_name, cache)
                if len(cache) % 10 == 0:
                    save_cache(cache_path, cache)
            except KeyboardInterrupt:
                raise
            except Exception as e:
                logging.error(f"Error processing {p}: {e}")
        save_cache(cache_path, cache)
        if not cfg.dry_run:
            build_index(cfg.output_dir)
    finally:
        # nothing special to close; llama.cpp cleans up via GC
        pass

    logging.info("Done.")


if __name__ == "__main__":
    main()
